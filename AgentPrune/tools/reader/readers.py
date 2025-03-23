#!/usr/bin/env python
# -*- coding: utf-8 -*-

from openai import OpenAI
import pdb

"""INSTALL
pip install openai --upgrade
pip install python-docx
pip install markdown
pip install PyPDF2
pip install openpyxl
pip install beautifulsoup4
pip install pylatexenc
pip install python-pptx
pip install xlrd
"""

import json
import os
import pandas as pd
import charset_normalizer
import docx
import markdown
import PyPDF2
import openpyxl
import yaml
import zipfile
import subprocess
from pathlib import Path
from abc import ABC, abstractmethod
from typing import Union, Any, Optional
from bs4 import BeautifulSoup
from pylatexenc.latex2text import LatexNodes2Text
from pptx import Presentation

from AgentPrune.utils.globals import Cost

from dotenv import load_dotenv
load_dotenv()
import aiohttp
import requests
from openai import OpenAI, AsyncOpenAI

OPENAI_API_KEY=os.getenv("OPENAI_API_KEY")


# Refs: https://platform.openai.com/docs/api-reference
# Refs: https://github.com/Significant-Gravitas/AutoGPT/blob/0e332c0c1221857f3ce96490f073c1c88bcbd367/autogpts/autogpt/autogpt/commands/file_operations_utils.py

class Reader(ABC):
    @abstractmethod
    def parse(self, file_path: Path) -> str:
        """ To be overriden by the descendant class """


class TXTReader(Reader):
    def parse(self, file_path: Path) -> str:
        content = charset_normalizer.from_path(file_path).best()
        return str(content)
    
class PDFReader(Reader):
    def parse(self, file_path: Path) -> str:
        content = PyPDF2.PdfReader(file_path)
        text = ""
        for page_idx in range(len(content.pages)):
            text += f'Page {page_idx + 1}\n' + content.pages[page_idx].extract_text()
        return text
    
class DOCXReader(Reader):
    def parse(self, file_path: Path) -> str:
        content = docx.Document(str(file_path))
        text = ""
        for i, para in enumerate(content.paragraphs):
            text += f'Page {i + 1}:\n' +  para.text
        return text

class JSONReader(Reader):
    def parse_file(file_path: Path) -> list:
        try:
            with open(file_path, "r") as f:
                data = json.load(f)
                #text = str(data)
            return data#text
        except:
            return []
    
    def parse(self, file_path: Path) -> str:
        try:
            with open(file_path, "r") as f:
                data = json.load(f)
                text = str(data)
            return text
        except:
            return ''
        
class JSONLReader(Reader):
    def parse_file(file_path: Path) -> list:
        with open(file_path, "r",encoding='utf-8') as f:
            lines = [json.loads(line) for line in f]
            #text = '\n'.join([str(line) for line in lines])
        return lines #text
    
    def parse(file_path: Path) -> str:
        with open(file_path, "r",encoding='utf-8') as f:
            lines = [json.loads(line) for line in f]
            text = '\n'.join([str(line) for line in lines])
        return text

class XMLReader(Reader):
    def parse(self, file_path: Path) -> str:
        with open(file_path, "r") as f:
            data = BeautifulSoup(f, "xml")
            text = data.get_text()
        return text

class YAMLReader(Reader):
    def parse(self, file_path: Path, return_str=True) -> Union[str, Any]:
        with open(file_path, "r") as f:
            data = yaml.load(f, Loader=yaml.FullLoader)
            text = str(data)
        if return_str:
            return text
        else:
            return data
    
class HTMLReader(Reader):
    def parse(self, file_path: Path) -> str:
        with open(file_path, "r") as f:
            data = BeautifulSoup(f, "html.parser")
            text = data.get_text()
        return text
    
class MarkdownReader(Reader):
    def parse(self, file_path: Path) -> str:
        with open(file_path, "r") as f:
            data = markdown.markdown(f.read())
            text = "".join(BeautifulSoup(data, "html.parser").findAll(string=True))
        return text

class LaTexReader(Reader):
    def parse(self, file_path: Path) -> str:
        with open(file_path, "r") as f:
            data = f.read()
        text = LatexNodes2Text().latex_to_text(data)
        return text



class AudioReader(Reader):
    @staticmethod
    def parse(file_path: Path) -> str:
        client = OpenAI(api_key=OPENAI_API_KEY)
        try:
            client = OpenAI()
            with open(file_path, "rb") as audio_file:
                transcript = client.audio.translations.create(
                    model="whisper-1",
                    file=audio_file
                )
            return transcript.text
        except Exception as e:
            return "Error transcribing audio file."

class PPTXReader(Reader): 
    def parse(self, file_path: Path) -> str:
        try:
            pres = Presentation(str(file_path))
            text = []
            for slide_idx, slide in enumerate(pres.slides):
                text.append(f"Slide {slide_idx + 1}:\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            return "Error reading PowerPoint file."

class ExcelReader(Reader):
    def parse(self, file_path: Path) -> str:
        try:
            excel_data = pd.read_excel(file_path, sheet_name=None)

            all_sheets_text = []
            for sheet_name, data in excel_data.items():
                all_sheets_text.append(f"Sheet Name: {sheet_name}\n{data.to_string()}\n")

            return "\n".join(all_sheets_text)
        except Exception as e:
            return "Error reading Excel file."

class XLSXReader(Reader):
    def parse(self, file_path: Path) -> str:
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        text = ""

        for sheet in workbook:
            text += f"\nSheet: {sheet.title}\n"
            for row in sheet.iter_rows(values_only=True):
                row_data = [str(cell) if cell is not None else "" for cell in row]
                text += "\t".join(row_data) + "\n"
        
        return text

class ZipReader(Reader):
    def parse(self, file_path: Path) -> str:
        #only support files that can be represented as text
        try:
            file_content = ""
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                extract_dir = file_path[:-4] + '/'
                zip_ref.extractall(extract_dir)
                reader = FileReader()
                for file_name in zip_ref.namelist():
                    file_content += f'File {file_name}:\n"{reader.read_file(extract_dir + file_name)}"\n'
            return file_content
        
        except zipfile.BadZipFile:
            print("Invalid ZIP file.")

        except Exception as e:
            print(f"Error reading ZIP file: {e}")


class PythonReader(Reader):
    def parse(self, file_path: Path) -> str:
        execution_result = ""
        error = ""
        file_content = ""
        try:
            completed_process = subprocess.run(["python", file_path], capture_output=True, text=True, check=True)
            execution_result = "Output:\n" + completed_process.stdout
        except subprocess.CalledProcessError as e:
            error = "Error:\n" + e.stderr
        except Exception as e:
            print(f"Error executing Python file: {e}")

        try:
            with open(file_path, "r") as file:
                file_content = "\nFile Content:\n" + file.read()
        except Exception as e:
            print(f"Error reading Python file: {e}")
        return file_content, execution_result, error


# Support 41 kinds of files.
READER_MAP = { 
    ".mp3": AudioReader(),
    ".m4a": AudioReader(),
    ".wav": AudioReader(),
    ".zip": ZipReader(),
    ".pptx": PPTXReader(),
    ".xlsx": ExcelReader(),
    ".xls": ExcelReader(),
    ".txt": TXTReader(),
    ".csv": TXTReader(),
    ".pdf": PDFReader(),
    ".docx": DOCXReader(),
    ".json": JSONReader(),
    ".jsonld": JSONReader(),
    ".jsonl": JSONLReader(),
    ".xml": XMLReader(),
    ".yaml": YAMLReader(),
    ".yml": YAMLReader(),
    ".html": HTMLReader(),
    ".htm": HTMLReader(),
    ".xhtml": HTMLReader(),
    ".md": MarkdownReader(),
    ".markdown": MarkdownReader(),
    ".tex": LaTexReader(),
    ".py": PythonReader(),
    ".pdb": TXTReader(),
}
    
class FileReader:
    def set_reader(self, suffix) -> None:
        self.reader = READER_MAP[suffix]
        print(f"Setting Reader to {type(self.reader).__name__}")

    def read_file(self, file_path: Path, task="describe the file") -> str:
        suffix = '.' + file_path.split(".")[-1]
        self.set_reader(suffix)
        file_content = self.reader.parse(file_path)
        print(f"Reading file {file_path} using {type(self.reader).__name__}")
        return file_content
    

class GeneralReader:
    def __init__(self):
        self.file_reader = FileReader()
        self.name = "General File Reader"
        self.description = """A general file reader support to formats: 'py', 'java', 'cpp', 'c', 'js', 
                              'css', 'html', 'htm', 'xml', 'txt', 'jsonl', 'csv', 'json', 
                              'jsonld', 'jsonl', 'yaml', 'yml', 'xlsx', 'xls', 'jpg', 'png', 
                              'jpeg', 'gif', 'bmp', 'mp3', 'wav', 'ogg', 'mp4', 'avi', 'mkv', 
                              'mov', 'pdf', 'doc', 'docx', 'ppt', 'pptx', 'md', 'markdown', 
                              'tex', 'zip', 'tar', 'gz', '7z', 'rar'.
                            """

    def read(self, task, file):

        files_content = ""
        file_content = self.file_reader.read_file(file, task)
        suffix = file.split(".")[-1]

        if suffix in ['py', 'java', 'cpp', 'c', 'js', 'css', 'html', 'htm', 'xml']:
            files_content += f'\nThe {suffix} file contains:\n---\n{file_content[0]}'
            if file_content[1] != '':
                files_content += f'\nExecution result:\n{file_content[1]}'
            if file_content[2] != '':
                files_content += f'\nExecution error message:\n{file_content[2]}'
            files_content += '\n---'

        elif suffix in ['txt', 'jsonl', 'csv', 'json', 'jsonld', 'jsonl', 'yaml', 'yml', 
                        'xlsx', 'xls', 'jpg', 'png', 'jpeg', 'gif', 'bmp', 'mp3', 'wav', 
                        'ogg', 'mp4', 'avi', 'mkv', 'mov', 'pdf', 'doc', 'docx', 'ppt', 
                        'pptx', 'md', 'markdown', 'tex', 'zip', 'tar', 'gz', '7z', 'rar']:
            files_content += f'\nThe {suffix} file contains:\n---\n{file_content}\n---'

        return files_content
    